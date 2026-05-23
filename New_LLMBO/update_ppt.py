#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
更新自我介绍PPT内容
"""
from pptx import Presentation
from pptx.util import Inches, Pt
import os

# PPT路径
ppt_path = 'D:/Users/aa133/Desktop/BO_Multi_12_20/New_LLMBO/自我介绍.pptx'
prs = Presentation(ppt_path)

def update_text_in_shape(shape, new_text):
    """更新shape中的文本"""
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.text = new_text
                return True
        if shape.text_frame.paragraphs:
            shape.text_frame.paragraphs[0].text = new_text
            return True
    return False

print("开始更新PPT内容...")
print("="*60)

# ==================== 第10页：研究经历 - 标题页 ====================
slide10 = prs.slides[9]
print("\n更新第10页：研究经历标题页")
for shape in slide10.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if "TITLE HERE" in text:
            update_text_in_shape(shape, "LLM增强多目标贝叶斯优化")
            print("  - 更新标题: LLM增强多目标贝叶斯优化")

# ==================== 第11页：研究背景与动机 ====================
slide11 = prs.slides[10]
print("\n更新第11页：研究背景与动机")
for shape in slide11.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if "TITLE HERE" in text:
            update_text_in_shape(shape, "电池快充协议优化问题")
            print("  - 更新标题: 电池快充协议优化问题")
        elif len(text) > 20:
            new_content = """• 问题定义：同时优化充电时间、峰值温升、容量衰减三个冲突目标
• 决策空间：5维参数（三阶段电流I1/I2/I3 + SOC宽度dSOC1/dSOC2）
• 核心挑战：传统BO纯数学驱动，缺乏电池领域知识引导"""
            update_text_in_shape(shape, new_content)
            print("  - 更新主要内容")

# ==================== 第12页：LLMBO方法框架 ====================
slide12 = prs.slides[11]
print("\n更新第12页：LLMBO方法框架")
for shape in slide12.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if "TITLE HERE" in text:
            update_text_in_shape(shape, "双触点的LLM增强策略")
            print("  - 更新副标题: 双触点的LLM增强策略")

# 更新四个要点的内容
content_items = [
    ("01", "WarmStart机制", "LLM生成高质量初始候选点，替代随机初始化，加速早期收敛"),
    ("02", "Region-Lifted GP", "LLM推荐有希望的区域，通过Mean Shift引导BO搜索方向"),
    ("03", "物理信息核", "LLM生成耦合矩阵，将电池物理嵌入GP核函数"),
    ("04", "自适应信任度", "根据HV提升动态调整对LLM建议的信任程度"),
]

idx = 0
for shape in slide12.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text in ["01", "02", "03", "04"] or (len(text) < 3 and text in ["01", "02", "03", "04"]):
            if idx < len(content_items):
                num, title, desc = content_items[idx]
                update_text_in_shape(shape, f"{num}\n{title}\n{desc}")
                print(f"  - 更新要点{num}: {title}")
                idx += 1

# ==================== 第13页：实验成果 ====================
slide13 = prs.slides[12]
print("\n更新第13页：实验成果")
for shape in slide13.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if "TITLE HERE" in text:
            update_text_in_shape(shape, "4-way对比验证LLMBO有效性")
            print("  - 更新副标题: 4-way对比验证LLMBO有效性")

# ==================== 第14页：学习探索 - 标题页 ====================
slide14 = prs.slides[13]
print("\n更新第14页：学习探索标题页")
for shape in slide14.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if "TITLE HERE" in text:
            update_text_in_shape(shape, "LLM技术学习笔记")
            print("  - 更新标题: LLM技术学习笔记")

# ==================== 第15页：PDF笔记展示 ====================
slide15 = prs.slides[14]
print("\n更新第15页：PDF笔记展示")
for shape in slide15.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if "TITLE HERE" in text:
            update_text_in_shape(shape, "系统学习大模型基础")
            print("  - 更新副标题: 系统学习大模型基础")

# 更新图片说明文本
image_captions = [
    "笔记1：注意力机制原理",
    "笔记2：位置编码与嵌入",
    "笔记3：多头注意力计算",
    "笔记4：前馈网络与残差连接",
]

cap_idx = 0
for shape in slide15.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if len(text) > 5 and cap_idx < len(image_captions):
            update_text_in_shape(shape, image_captions[cap_idx])
            print(f"  - 更新图片说明{cap_idx+1}: {image_captions[cap_idx]}")
            cap_idx += 1

# ==================== 第16页：学习总结 ====================
slide16 = prs.slides[15]
print("\n更新第16页：学习总结")
for shape in slide16.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if "TITLE HERE" in text:
            update_text_in_shape(shape, "理论基础与实践结合")
            print("  - 更新副标题: 理论基础与实践结合")

# ==================== 第17页：展望 - 标题页 ====================
slide17 = prs.slides[16]
print("\n更新第17页：展望标题页")
for shape in slide17.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if "TITLE HERE" in text:
            update_text_in_shape(shape, "多模态大模型探索")
            print("  - 更新副标题: 多模态大模型探索")

# 更新要点内容
future_points = [
    "多模态理解：学习CLIP、LLaVA等模型，理解图文对齐机制",
    "生成式AI：探索扩散模型在可控生成中的应用",
    "智能体交互：了解GUI Agent基础，关注人机交互新范式",
    "持续学习：跟进LLM前沿进展，培养科研敏感度",
]

pt_idx = 0
for shape in slide17.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if len(text) > 30 and pt_idx < len(future_points):
            update_text_in_shape(shape, future_points[pt_idx])
            print(f"  - 更新要点{pt_idx+1}: {future_points[pt_idx][:30]}...")
            pt_idx += 1

# ==================== 第18页：多模态兴趣 ====================
slide18 = prs.slides[17]
print("\n更新第18页：多模态兴趣")
for shape in slide18.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if "TITLE HERE" in text:
            update_text_in_shape(shape, "与张驰老师方向的契合点")
            print("  - 更新副标题: 与张驰老师方向的契合点")
        elif len(text) > 20:
            new_text = """• 我已初步了解CLIP的视觉-语言对齐机制和LLaVA的多模态对话能力
• 在LLMBO项目中积累的LLM调优经验可直接迁移到多模态场景
• 对GUI Agent方向有浓厚兴趣，希望探索模型与真实界面的交互
• 愿意从基础做起，在导师指导下逐步深入多模态大模型研究"""
            update_text_in_shape(shape, new_text)
            print("  - 更新主要内容")

# ==================== 第19页：具体研究兴趣 ====================
slide19 = prs.slides[18]
print("\n更新第19页：具体研究兴趣")
for shape in slide19.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if "TITLE HERE" in text:
            update_text_in_shape(shape, "希望在博士阶段深入探索")
            print("  - 更新副标题: 希望在博士阶段深入探索")

# 更新六个兴趣点
interest_points = [
    "视觉-语言预训练",
    "多模态指令微调",
    "视觉Agent基础",
    "跨模态表示学习",
    "生成式多模态模型",
    "高效多模态架构",
]

int_idx = 0
for shape in slide19.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if len(text) > 5 and int_idx < len(interest_points):
            update_text_in_shape(shape, interest_points[int_idx])
            print(f"  - 更新兴趣点{int_idx+1}: {interest_points[int_idx]}")
            int_idx += 1

# ==================== 第20页：为什么选择西湖 ====================
slide20 = prs.slides[19]
print("\n更新第20页：为什么选择西湖")
for shape in slide20.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if "TITLE HERE" in text:
            update_text_in_shape(shape, "博士申请动机")
            print("  - 更新副标题: 博士申请动机")

# ==================== 第21页：总结 ====================
slide21 = prs.slides[20]
print("\n更新第21页：总结")
for shape in slide21.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if "TITLE HERE" in text:
            update_text_in_shape(shape, "个人总结")
            print("  - 更新标题: 个人总结")

# ==================== 第22页：致谢 ====================
slide22 = prs.slides[21]
print("\n更新第22页：致谢")
for shape in slide22.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if len(text) < 10 and ("汇报" in text or "感谢" in text or len(text) < 5):
            update_text_in_shape(shape, "感谢聆听\n恳请批评指正")
            print("  - 更新致谢语")

# 保存修改后的PPT
output_path = 'D:/Users/aa133/Desktop/BO_Multi_12_20/New_LLMBO/自我介绍_已更新.pptx'
prs.save(output_path)

print("\n" + "="*60)
print(f"PPT更新完成！")
print(f"保存路径: {output_path}")
print("="*60)
